#!/usr/bin/env python3
"""Generate dev_summary.pptx — compact 16-slide version with SVG-rendered diagrams."""
import io
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import cairosvg

# 复用 build_html.py 里手工调好的 SVG 图（架构 / 协同四层次 / L0-L4 / 卡死防御 等）
sys.path.insert(0, str(Path(__file__).resolve().parent))
from build_html import (  # noqa: E402
    svg_architecture,
    svg_collab_levels,
    svg_4layer_defense,
    svg_harness_l0_l4,
    svg_implementation_results,
)

# Colors
# 主配色：华为红（品牌色 #C7000B）— 仅用于"重点突出"（封面 / 章节标题 / callout）
PRIMARY = RGBColor(0xC7, 0x00, 0x0B)
DARK = RGBColor(0x20, 0x20, 0x20)
GRAY = RGBColor(0x60, 0x60, 0x60)
# 浅色衬底改为低饱和红粉色（# FFF1F2），与 PRIMARY 同色系协调
LIGHT_BG = RGBColor(0xFF, 0xF1, 0xF2)
# 表格标题 / 非重点区分块的浅灰底（替换原来的红 / 绿底）
LIGHT_HEADER = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_HEADER_BORDER = RGBColor(0xD0, 0xD0, 0xD0)
RED = RGBColor(0xE5, 0x39, 0x35)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xF5, 0x7C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_CN = "Microsoft YaHei"
FONT_MONO = "Consolas"

# Font sizes per spec
TITLE_PT = 24
BODY_LG = 16
BODY_MD = 14
BODY_SM = 12

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK_LAYOUT = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


def add_textbox(slide, left, top, width, height, text, *,
                font_size=BODY_MD, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, font_name=FONT_CN):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_header(slide, title_text):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.5),
                title_text, font_size=TITLE_PT, bold=True, color=PRIMARY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.5), Inches(0.85),
                                    Inches(12.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()


def add_page_number(slide, num, total):
    add_textbox(slide, Inches(11.8), Inches(7.05), Inches(1.3), Inches(0.3),
                f"{num} / {total}", font_size=10, color=GRAY,
                align=PP_ALIGN.RIGHT)


def add_table(slide, left, top, width, height, headers, rows,
              header_color=LIGHT_HEADER, alt_color=LIGHT_BG,
              font_size=BODY_SM, header_font_size=BODY_SM,
              col_widths=None, header_text_color=DARK):
    """表格。默认 header 浅灰底 + 深字（用户要求："非重点突出部分用浅灰底"）；
    重点强调的表格可显式传 header_color=PRIMARY + header_text_color=WHITE。"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.margin_left = Emu(40000)
        tf.margin_right = Emu(40000)
        tf.margin_top = Emu(15000)
        tf.margin_bottom = Emu(15000)
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(h)
        run.font.name = FONT_CN
        run.font.size = Pt(header_font_size)
        run.font.bold = True
        run.font.color.rgb = header_text_color

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.margin_left = Emu(40000)
            tf.margin_right = Emu(40000)
            tf.margin_top = Emu(10000)
            tf.margin_bottom = Emu(10000)
            tf.word_wrap = True
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT_CN
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK
    return table


def add_code_block(slide, left, top, width, height, code, *, font_size=BODY_SM):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF4, 0xF4, 0xF4)
    box.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(120000)
    tf.margin_right = Emu(120000)
    tf.margin_top = Emu(80000)
    tf.margin_bottom = Emu(80000)
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_MONO
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK


def embed_svg(slide, svg_str, left, top, width, *, render_width_px=2000):
    """把 SVG 字符串渲染成 PNG 后插入幻灯片。复用 build_html.py 的 SVG 函数，
    保证 PPT 与 HTML 渲染版的视觉一致。"""
    png_bytes = cairosvg.svg2png(
        bytestring=svg_str.encode("utf-8"),
        output_width=render_width_px,
    )
    return slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=width)


def add_callout(slide, left, top, width, height, text, *,
                bg=LIGHT_BG, border=PRIMARY, font_size=BODY_MD,
                color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = border
    if not text:
        return
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(150000)
    tf.margin_right = Emu(150000)
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CN
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


# =================================================================
# Slide 1: Cover + Agenda
# =================================================================
s = add_slide()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), SLIDE_H)
accent.fill.solid()
accent.fill.fore_color.rgb = PRIMARY
accent.line.fill.background()

add_textbox(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.6),
            "AI 辅助联网游戏开发", font_size=BODY_LG, color=GRAY)
add_textbox(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.0),
            "完整实践总结", font_size=44, bold=True, color=PRIMARY)
add_textbox(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.5),
            "沟通牌  ×  Claude Code  ×  54 PR · 有效开发 18 天",
            font_size=BODY_LG, color=DARK)

add_textbox(s, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.4),
            "本次内容（约 20 分钟）",
            font_size=BODY_MD, bold=True, color=PRIMARY)
items = [
    "①  项目背景与开发全貌    ②  架构设计与遗憾",
    "③  问题全景：人工 vs AI    ④  人机协同模式",
    "⑤  AI 质量改进路径：多 Claude 协同 + 跨 vendor",
    "⑥  关键技术修复    ⑦  经验总结与后续行动",
]
y = 4.5
for it in items:
    add_textbox(s, Inches(1.0), Inches(y), Inches(11.0), Inches(0.4),
                it, font_size=BODY_MD, color=DARK)
    y += 0.5

# =================================================================
# Slide 2: Background + Tech Stack + Code Volume + Timeline
# =================================================================
s = add_slide()
add_header(s, "一、项目背景与开发全貌")

add_textbox(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
            "沟通牌：4 副牌 216 张，6 人 3v3 卡牌游戏（单机 + 联网 + Web 三模式）",
            font_size=BODY_MD, color=DARK)

add_textbox(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.4),
            "技术栈", font_size=BODY_LG, bold=True, color=PRIMARY)
tech = [
    ["共享逻辑", "Kotlin Multiplatform (android+jvm+wasmJs)"],
    ["Android", "Kotlin + Coroutines + Flow + OkHttp + XML"],
    ["Web", "Compose Multiplatform 1.6.10 / Wasm-JS"],
    ["服务端", "Ktor 2.3.6 + Netty + :shared 依赖"],
    ["反代部署", "Caddy 80/443 → :8080 + systemd auto-deploy"],
]
add_table(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(2.3),
            ["层", "技术"], tech, font_size=BODY_SM,
            col_widths=[Inches(1.5), Inches(4.5)])

add_textbox(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(0.4),
            "代码规模 (PR #54 后)  ·  约 19,250 行",
            font_size=BODY_LG, bold=True, color=PRIMARY)
vol = [
    [":apps:android (UI+网络)", "19 文件 · ~6,300 行"],
    [":apps:web (Compose MP)", "23 文件 · ~3,470 行"],
    [":shared (KMP commonMain)", "9 文件 · ~2,670 行"],
    [":server (Ktor)", "4 文件 · ~1,960 行"],
    ["测试 + Android XML", "24 文件 · ~4,850 行"],
]
add_table(s, Inches(7.0), Inches(2.0), Inches(6.0), Inches(2.3),
            ["模块", "规模"], vol, font_size=BODY_SM,
            col_widths=[Inches(2.4), Inches(3.6)])

add_textbox(s, Inches(0.5), Inches(4.45), Inches(12.3), Inches(0.4),
            "12 阶段时间线  ·  54 PR  ·  ~170 commit  ·  有效开发 18 天",
            font_size=BODY_LG, bold=True, color=PRIMARY)
phases = [
    ["2026-02-02 / 07~12 / 24", "单机游戏 (#1–14)", "引擎/牌型/AI/结算 · 11 轮人工 UI 反馈"],
    ["2026-04-30", "联网首版 (#16)", "服务端 + 网络层 · 一次性 6,529 行"],
    ["2026-05-01~03", "部署 (#17–31)", "CI/Gradle/cleartext/503/Lobby UI"],
    ["2026-05-07", "深度修复 (#34)", "AI 全量审查×4轮 / 8 commit / ~50 Bug"],
    ["2026-05-08", "KMP 重构 + Web (#35)", "抽 :shared / Compose MP wasmJs"],
    ["2026-05-08", "Harness H1–H5 (#36–43)", "hooks / TDD-gate / pr-reviewer"],
    ["2026-05-09", "Web UI 4 阶段 (#47–50)", "菜单/响应式/视觉/Android 同等"],
    ["2026-05-10", "AI 接管 + 单机修复 (#51–54)", "G34-G38 · PROTOCOL_VERSION 3"],
]
add_table(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.0),
            ["时间", "阶段", "主要内容"], phases, font_size=BODY_SM,
            col_widths=[Inches(2.6), Inches(3.7), Inches(6.0)])

# =================================================================
# Slide 3: Architecture (diagram + decisions + regrets)
# =================================================================
s = add_slide()
add_header(s, "二、架构设计：多端共享 + 服务端权威")

# 用 build_html.py 中的整体架构 SVG 替代 ASCII（视觉与 dev_summary.html 一致）
embed_svg(s, svg_architecture(), Inches(0.4), Inches(1.0), width=Inches(6.5))

add_textbox(s, Inches(7.0), Inches(1.05), Inches(6.0), Inches(0.4),
            "关键架构决策", font_size=BODY_LG, bold=True, color=PRIMARY)
decisions = [
    ["状态同步", "全量状态 + 单调 version"],
    ["并发模型", "每房间 Mutex（修改在锁内）"],
    ["重连机制", "sessionToken=playerId · 30s 内恢复"],
    ["AI 替补", "isAISubstitute 标记，不删玩家槽"],
    ["共享逻辑", ":shared KMP（编译期保证一致）"],
    ["协议版本", "PROTOCOL_VERSION=3 · 握手时拒老客户端"],
]
add_table(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(2.8),
            ["决策", "选择"], decisions, font_size=BODY_SM,
            col_widths=[Inches(1.5), Inches(4.5)])

add_textbox(s, Inches(7.0), Inches(4.5), Inches(6.0), Inches(0.4),
            "架构演进：从遗憾到修复", font_size=BODY_LG, bold=True, color=GREEN)
evolution = [
    ["KMP 共享模块", "✅ PR #35 + H3：编译期唯一份"],
    ["协议版本号", "✅ PR-H3：PROTOCOL_VERSION + 握手"],
    ["事件溯源", "⚪ 未规划，全量同步够用"],
]
# 表格 header 用浅灰底（与全局规则一致；不再用 GREEN 填底）
add_table(s, Inches(7.0), Inches(4.95), Inches(6.0), Inches(1.5),
            ["遗憾", "状态"], evolution, font_size=BODY_SM,
            col_widths=[Inches(2.0), Inches(4.0)])

add_callout(s, Inches(0.5), Inches(6.6), Inches(12.5), Inches(0.55),
            "核心原则：服务端权威  ·  客户端乐观响应  ·  全量状态 + version 同步",
            bg=PRIMARY, border=PRIMARY,
            font_size=BODY_MD, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# =================================================================
# Slide 4: Problem Discovery Overview
# =================================================================
s = add_slide()
add_header(s, "三、问题全景：人工 vs AI")

y = 1.1
cards = [
    ("🔴 人工测试 / 反馈", "~35", "27%",
     "UI 体验 · 部署环境 · 运行时崩溃", "真机发现", RED),
    ("🔵 Claude Code 主会话", "~55", "43%",
     "全量扫描 · 跨文件链路 · 并发陷阱", "opus-4-7 / sonnet-4-6", PRIMARY),
    ("🟢 Claude pr-reviewer", "~15", "12%",
     "独立 context · 功能完整 · 协议契约", "PR-H5 后引入", GREEN),
    ("🟡 ChatGPT Codex Bot", "~13", "10%",
     "语句级边界 · entropy · UI 文案", "chatgpt-codex-connector", ORANGE),
]
x_positions = [0.5, 3.7, 6.9, 10.1]
for (title, num, pct, desc, model, color), x in zip(cards, x_positions):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y),
                                Inches(3.0), Inches(2.6))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    add_textbox(s, Inches(x), Inches(y + 0.1), Inches(3.0), Inches(0.4),
                title, font_size=BODY_SM, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(x), Inches(y + 0.55), Inches(3.0), Inches(0.85),
                num, font_size=36, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(x), Inches(y + 1.45), Inches(3.0), Inches(0.3),
                pct, font_size=BODY_SM, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(x), Inches(y + 1.8), Inches(3.0), Inches(0.4),
                desc, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(x), Inches(y + 2.2), Inches(3.0), Inches(0.3),
                model, font_size=9, color=WHITE,
                align=PP_ALIGN.CENTER)

add_callout(s, Inches(0.5), Inches(4.0), Inches(12.4), Inches(3.0),
            "", bg=LIGHT_BG, border=PRIMARY)
add_textbox(s, Inches(0.8), Inches(4.15), Inches(12.0), Inches(0.5),
            "💡 核心发现：4 种视角几乎不重叠 → 总计 ~128 个问题",
            font_size=BODY_LG, bold=True, color=PRIMARY)
add_textbox(s, Inches(0.8), Inches(4.75), Inches(12.0), Inches(0.4),
            "·  人工：截图 / 真机崩溃 / 部署环境（动态运行时问题）",
            font_size=BODY_SM, color=DARK)
add_textbox(s, Inches(0.8), Inches(5.15), Inches(12.0), Inches(0.4),
            "·  Claude 主会话：跨文件链路 / 并发陷阱 / 工具链兼容（静态全量扫描）",
            font_size=BODY_SM, color=DARK)
add_textbox(s, Inches(0.8), Inches(5.55), Inches(12.0), Inches(0.4),
            "·  Claude pr-reviewer：独立 context · 功能完整性 · 跨文件契约（PR #53 P1 #1：fabricated symbol）",
            font_size=BODY_SM, color=DARK)
add_textbox(s, Inches(0.8), Inches(5.95), Inches(12.0), Inches(0.4),
            "·  Codex bot：UUID 截断 / loading 卡死 / processAITurn race（语句级边界）",
            font_size=BODY_SM, color=DARK)
add_textbox(s, Inches(0.8), Inches(6.45), Inches(12.0), Inches(0.4),
            "→  4 视角缺一不可：盲区互补，单独跑一个至少漏一类问题",
            font_size=BODY_MD, bold=True, color=PRIMARY)

# =================================================================
# Slide 5: Human-discovered Issues (3 phases compacted)
# =================================================================
s = add_slide()
add_header(s, "四、人工发现的问题（~32 个，3 阶段）")

add_textbox(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
            "单机阶段（11 个）  ·  反复 UI 打磨",
            font_size=BODY_MD, bold=True, color=PRIMARY)
sp = [
    ["UI 显示", "看不到玩家出牌 · 玩家ID映射错位 · 字体不统一 · 布局拥挤 · 卡片圆角丢失"],
    ["游戏逻辑", "队伍积分只显示个人 · 手牌顺序乱 · AI 滥用炸弹"],
    ["环境兼容", "炸弹重叠比例 · APK 签名不匹配 · 五子棋图标旧 Android 崩溃"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
            ["类型", "症状（合并相似项）"], sp, font_size=BODY_SM,
            col_widths=[Inches(2.0), Inches(10.3)])

add_textbox(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.4),
            "部署阶段（9 个）  ·  环境与构建",
            font_size=BODY_MD, bold=True, color=PRIMARY)
dep = [
    ["构建编译", "CI 拉入服务端模块 · Gradle Wrapper 缺失 · 联网代码 3 处编译错误 · JVM 工具链 · 视图 ID"],
    ["部署连通", "服务器 URL · Android 9+ cleartext · 503 错误"],
    ["UI 表述", "「开始游戏」按钮让用户困惑"],
]
add_table(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(1.5),
            ["类型", "症状（合并相似项）"], dep, font_size=BODY_SM,
            col_widths=[Inches(2.0), Inches(10.3)])

add_textbox(s, Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.4),
            "联网游戏阶段（12 个）  ·  ⭐ 核心痛点",
            font_size=BODY_MD, bold=True, color=RED)
mp = [
    ["UI / 大厅", "loading 永久卡住 · 昵称限制 · 大厅崩溃 · AI 离线 · 无房间列表 · 无踢人"],
    ["📷 游戏卡死", "等待电脑 54 出牌（反复 3 次复现，触发 4 轮 AI 全量自查）"],
    ["📷 分数错误", "已收分全是 0"],
]
add_table(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.5),
            ["类型", "症状（合并相似项）"], mp, font_size=BODY_SM,
            col_widths=[Inches(2.0), Inches(10.3)])

# =================================================================
# Slide 6: AI-discovered Issues (Claude + Codex)
# =================================================================
s = add_slide()
add_header(s, "四、AI 发现的问题（~38 个）")

# Section 1: Claude Code Agent
add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
            "🔵 AI #1：Claude Code Agent（~35 个，4 轮审查）",
            font_size=BODY_MD, bold=True, color=PRIMARY)
add_textbox(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
            "模型：claude-opus-4-7（1M 上下文）+ claude-sonnet-4-6  ·  工作模式：Level 4「自查自纠」",
            font_size=10, color=GRAY)

claude_rounds = [
    ["第 1 轮 综合审查 (~20)",
     "协议/序列化 4 · 会话/重连 3 · 房间状态机 4 · UI/状态同步 6 · 其他 3"],
    ["第 2-3 轮 深层 (~8)",
     "回合错位 · ArrayList 并发 · AI 回退链缺失 · collectedScore 硬编码 0"],
    ["第 4 轮 根因 (~7) ⭐",
     "WebSocket CONNECTING send() 静默失败 · 多协程无锁并发 · 结算公式漏算"],
]
add_table(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(1.95),
            ["审查轮次", "主要问题"], claude_rounds, font_size=BODY_SM,
            col_widths=[Inches(3.0), Inches(9.3)])

# Section 2: ChatGPT Codex Review Bot
add_textbox(s, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.4),
            "🟡 AI #2：ChatGPT Codex Review Bot（3 个）",
            font_size=BODY_MD, bold=True, color=ORANGE)
add_textbox(s, Inches(0.5), Inches(4.45), Inches(12.3), Inches(0.4),
            "Agent：chatgpt-codex-connector[bot]  ·  PR 自动审查触发  ·  按 P1/P2 标注优先级",
            font_size=10, color=GRAY)

codex_findings = [
    ["#29", "P1", "Application.kt:48", "UUID 截断到 8 字符 → 碰撞致跨用户混乱", "❌→✅ 本次修复"],
    ["#31", "P1", "LobbyActivity.kt:219", "重连后 loading 遮罩可能永久卡住", "✅ PR #33 已修复"],
    ["#33", "P2", "MainActivity.kt", "UI 提示按房间名加入但服务端不支持", "✅ UI 重构后修复"],
]
add_table(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(1.6),
            ["PR", "优先级", "位置", "意见", "处置"], codex_findings,
            font_size=BODY_SM,
            col_widths=[Inches(0.8), Inches(0.9), Inches(2.4),
                        Inches(5.5), Inches(2.7)])

# Bottom callout
add_callout(s, Inches(0.5), Inches(6.7), Inches(12.4), Inches(0.55),
            "💡 关键观察：Claude 偏向「全局逻辑链路」 vs Codex 偏向「细粒度风险点」 — 两 AI 几乎不重叠",
            bg=LIGHT_BG, border=PRIMARY,
            font_size=BODY_SM, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# =================================================================
# Slide 7: Collaboration — 4 levels + Matrix
# =================================================================
s = add_slide()
add_header(s, "五、人机协同：4 层次 + 实际分工")

add_textbox(s, Inches(0.5), Inches(1.05), Inches(6.5), Inches(0.4),
            "协同的 4 个层次", font_size=BODY_LG, bold=True, color=PRIMARY)
# 复用 build_html.py 的 SVG（与渲染版 HTML 视觉一致）
embed_svg(s, svg_collab_levels(), Inches(0.4), Inches(1.5), width=Inches(6.6))

add_textbox(s, Inches(7.2), Inches(1.05), Inches(6.0), Inches(0.4),
            "实际任务分工矩阵", font_size=BODY_LG, bold=True, color=PRIMARY)
matrix = [
    ["需求定义", "100%", "—"],
    ["架构设计", "70%", "30%"],
    ["编码实现", "5%", "95%"],
    ["UI 调试", "60%", "40%"],
    ["逻辑 Bug 排查", "20%", "80%"],
    ["部署/网络", "80%", "20%"],
    ["文档编写", "10%", "90%"],
    ["代码审查", "30%", "70%"],
]
add_table(s, Inches(7.2), Inches(1.5), Inches(5.8), Inches(4.6),
            ["任务", "人工", "AI"], matrix, font_size=BODY_SM,
            col_widths=[Inches(2.4), Inches(1.7), Inches(1.7)])

add_callout(s, Inches(0.5), Inches(6.4), Inches(12.4), Inches(0.6),
            "本项目 L3 + L4 占用约 70% 协同时间  ·  开放性指令激发全量审查",
            bg=LIGHT_BG, border=PRIMARY,
            font_size=BODY_MD, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# =================================================================
# Slide 8: Strengths + Anti-patterns + Best practices + Efficiency
# =================================================================
s = add_slide()
add_header(s, "五、优势分工 · 反模式 · 效率")

add_textbox(s, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.4),
            "✅ AI 显著优于人工", font_size=BODY_MD, bold=True, color=GREEN)
ai_pts = [
    "·  全量代码审查（35 问题 / 90 文件）",
    "·  跨文件链路追踪",
    "·  重复模式识别（5 处一次发现）",
    "·  测试用例生成（15 边界用例）",
]
y = 1.5
for pt in ai_pts:
    add_textbox(s, Inches(0.6), Inches(y), Inches(6.0), Inches(0.35),
                pt, font_size=BODY_SM, color=DARK)
    y += 0.32

add_textbox(s, Inches(6.8), Inches(1.05), Inches(6.0), Inches(0.4),
            "❌ 人工不可替代", font_size=BODY_MD, bold=True, color=RED)
hu_pts = [
    "·  真机环境验证（cleartext / 503）",
    "·  时序竞争复现（网络抖动 / 并发）",
    "·  用户体验判断（视觉感知）",
    "·  部署 & 业务决策",
]
y = 1.5
for pt in hu_pts:
    add_textbox(s, Inches(6.9), Inches(y), Inches(6.0), Inches(0.35),
                pt, font_size=BODY_SM, color=DARK)
    y += 0.32

add_textbox(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.4),
            "🚨 协同反模式（要避免）",
            font_size=BODY_MD, bold=True, color=RED)
anti = [
    ["过度信任", "AI 说「已修复」就直接合入 → 反复复现"],
    ["模糊指令", "「把这个 bug 修了」 → 只修表象"],
    ["一次到位幻想", "期待一次解决所有问题 → 本项目经历 4 轮"],
    ["跳过验证", "AI 修完直接发布 → 真机问题永远暴露不出来"],
]
add_table(s, Inches(0.5), Inches(3.5), Inches(7.5), Inches(2.4),
            ["反模式", "后果"], anti, font_size=BODY_SM,
            col_widths=[Inches(2.0), Inches(5.5)])

add_textbox(s, Inches(8.3), Inches(3.0), Inches(4.7), Inches(0.4),
            "📊 效率数据",
            font_size=BODY_MD, bold=True, color=PRIMARY)
eff = [
    ["人工总投入", "~30 小时"],
    ["AI 等效工作量", "~300 小时"],
    ["提速比", "约 10 倍"],
    ["单次成功率", "~12%"],
]
add_table(s, Inches(8.3), Inches(3.5), Inches(4.7), Inches(2.4),
            ["指标", "数值"], eff, font_size=BODY_SM,
            col_widths=[Inches(2.5), Inches(2.2)])

add_textbox(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.4),
            "💡 高效协同 6 条实践",
            font_size=BODY_MD, bold=True, color=PRIMARY)
add_textbox(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.4),
            "①症状要具体  ②截图优于文字  ③允许多轮迭代  ④关键决策人工拍板  ⑤人工守发布闸门  ⑥开放性指令激发全量审查",
            font_size=BODY_SM, color=DARK)

# =================================================================
# Slide 9: Reflection — Why Claude code has bugs (root cause)
# =================================================================
s = add_slide()
add_header(s, "六、反思：为什么 AI 写的代码 Bug 不少？")

add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
            "🔵 生成阶段必然有 Bug 的 5 个结构性原因",
            font_size=BODY_MD, bold=True, color=PRIMARY)
gen_reasons = [
    ["1", "生成 vs 验证是不同认知任务", "canBeat 套用「同张数比大小」常用模式，忽略大炸弹直胜"],
    ["2", "生成时无执行反馈，时序/并发是盲区", "WebSocket CONNECTING 时 send() 静默失败"],
    ["3", "自然语言规格隐式不完整", "「重连后游戏继续」未明确秒数和保留状态"],
    ["4", "模式匹配编码「常用 ≠ 正确」", "UUID.take(8) / ArrayList — 常用但联网场景错"],
    ["5", "跨文件一致性是结构盲区", "单机 CardRules vs 服务端 canBeat — 双份 3 处不一致"],
]
add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(2.5),
            ["#", "原因", "本项目例子"], gen_reasons, font_size=BODY_SM,
            col_widths=[Inches(0.5), Inches(4.5), Inches(7.3)])

add_textbox(s, Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.4),
            "🔁 多轮自审仍能发现新问题的 4 个原因",
            font_size=BODY_MD, bold=True, color=ORANGE)
multi_reasons = [
    ["1", "每轮在问不同的问题", "R1:「代码整洁吗？」 R4:「为什么 3 次修了还卡？」"],
    ["2", "自审有确认偏差，用户「还卡住」才打破", "4 轮逐层挖到 send() 静默失败"],
    ["3", "症状被消除后下一层根因才暴露", "canBeat → 回退链 → Mutex → 兜底 必须按序解锁"],
    ["4", "注意力有限，单轮无法全部仔细看", "1M 上下文 ≠ 全部能仔细分析"],
]
add_table(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.0),
            ["#", "原因", "本项目体现"], multi_reasons, font_size=BODY_SM,
            col_widths=[Inches(0.5), Inches(4.5), Inches(7.3)])

add_callout(s, Inches(0.5), Inches(6.75), Inches(12.4), Inches(0.55),
            "💡 启示：AI 写代码不是「一次到位」，而是「快速迭代到位」 — 接受多轮，但用工具/流程压低轮次",
            bg=LIGHT_BG, border=PRIMARY,
            font_size=BODY_SM, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# =================================================================
# Slide 10: 5 Multi-Claude Collaboration Patterns
# =================================================================
s = add_slide()
add_header(s, "六、多 Claude 模型协同的 5 种模式")

# Top: available models
add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
            "可用模型",
            font_size=BODY_MD, bold=True, color=PRIMARY)
models = [
    ["Claude Opus 4.7（1M）", "推理深度最强，全局视角", "架构 / 根因分析 / 深度审查"],
    ["Claude Sonnet 4.6", "平衡速度与能力", "主要实现 / PR review"],
    ["Claude Haiku 4.5", "极快、便宜", "静态扫描 / 测试生成 / 批量检查"],
]
add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.6),
            ["模型", "特点", "适合的角色"], models, font_size=BODY_SM,
            col_widths=[Inches(3.0), Inches(4.5), Inches(4.8)])

# Bottom: 5 patterns
add_textbox(s, Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.4),
            "5 种协同模式（按可行性排序）",
            font_size=BODY_MD, bold=True, color=PRIMARY)
patterns = [
    ["①", "开新会话 = 等价的「另一个 Claude」",
     "上下文重置打破自审偏差，零成本，本项目 4 轮自查实质即此 — 推荐"],
    ["②", "Generator / Reviewer 分工",
     "Opus 架构 → Sonnet 实现 → Haiku 静态扫描 → Opus（新会话）根因审查"],
    ["③", "对抗式审查（Adversarial）",
     "A 实现 / B 攻击「找崩溃场景」 / C 仲裁 — 适合金钱/分数/安全代码"],
    ["④", "TDD 反向流",
     "Haiku 先写边界用例 → Sonnet 实现到通过 → Opus 审查覆盖率 — 最压低 Bug"],
    ["⑤", "Self-Consistency 校验",
     "同一任务 Opus 跑 3 次对比，差异处标记为「不确定区」 — 关键代码"],
]
add_table(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(3.4),
            ["#", "模式", "做法 / 适用场景"], patterns, font_size=BODY_SM,
            col_widths=[Inches(0.5), Inches(3.5), Inches(8.3)])

# =================================================================
# Slide 11: Coverage ceiling + Cross-vendor + Recommended workflow
# =================================================================
s = add_slide()
add_header(s, "六、协同天花板 & 推荐工作流")

# Top: ceiling
add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
            "⚠ 多 Claude 协同的天花板：相关性盲区",
            font_size=BODY_MD, bold=True, color=RED)
add_textbox(s, Inches(0.7), Inches(1.45), Inches(12.0), Inches(0.4),
            "原因：所有 Claude 共享同一训练语料 + 同一训练目标 + 同一架构",
            font_size=BODY_SM, color=DARK)
add_textbox(s, Inches(0.7), Inches(1.85), Inches(12.0), Inches(0.4),
            "→ Codex Bot 找到的 3 个问题（UUID 截断 / loading 卡死 / UI 不一致）多 Claude 也大概率找不出",
            font_size=BODY_SM, color=GRAY)

# Middle: Cross-vendor strategy
add_textbox(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.4),
            "📊 真正补盲区的策略对比",
            font_size=BODY_MD, bold=True, color=PRIMARY)
cross = [
    ["多个 Claude 实例", "全局架构 + 逻辑链路（多角度）", "中等"],
    ["Claude + Codex（OpenAI）", "增加细粒度风险点", "较高"],
    ["Claude + Codex + Gemini", "不同训练目标，覆盖最广", "最高"],
    ["Claude + Detekt / SpotBugs / kover", "规则化盲区", "必要"],
]
add_table(s, Inches(0.5), Inches(2.85), Inches(12.3), Inches(1.8),
            ["组合", "找到的问题类型", "互补程度"], cross, font_size=BODY_SM,
            col_widths=[Inches(4.5), Inches(5.5), Inches(2.3)])

# Bottom: Recommended workflow
add_textbox(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.4),
            "🚀 给本项目的推荐工作流（如果重做联网模块）",
            font_size=BODY_MD, bold=True, color=GREEN)

workflow = """开发：  Opus 4.7 (1M)  架构 + 协议    Sonnet 4.6  实现    Haiku 4.5  静态扫描

PR 4 关：  Claude PR Review (新会话)  +  Codex Bot  +  Detekt  +  人工真机验证

测试驱动：  关键路径（结算 / 协议 / 并发）必须 TDD，Haiku 4.5 批量生成边界用例"""
add_code_block(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.6),
                workflow, font_size=BODY_SM)

add_callout(s, Inches(0.5), Inches(7.0), Inches(12.4), Inches(0.4),
            "💡 单 Claude 多轮 = 时间换覆盖率  ·  多 Claude = 视角换覆盖率  ·  多 vendor + 静态 + 真机 = 异构换覆盖率（边际收益最大）",
            bg=LIGHT_BG, border=PRIMARY,
            font_size=10, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# =================================================================
# Slide 12: Key Technical Fixes (6 in 1, compacted)
# =================================================================
s = add_slide()
add_header(s, "七、关键技术修复（6 处典型）")

# 左半：3 个游戏逻辑层修复
add_textbox(s, Inches(0.5), Inches(1.05), Inches(6.2), Inches(0.4),
            "🎮 游戏逻辑层", font_size=BODY_LG, bold=True, color=PRIMARY)

add_textbox(s, Inches(0.5), Inches(1.5), Inches(6.2), Inches(0.35),
            "1️⃣ 游戏卡死 — 四层防御",
            font_size=BODY_SM, bold=True, color=DARK)
fix1 = """L1 canBeat 炸弹大张数胜
L2 AI 三级回退
L3 每房间 Mutex
L4 broadcastForceAdvance 兜底"""
add_code_block(s, Inches(0.5), Inches(1.85), Inches(6.2), Inches(1.2),
                fix1, font_size=10)

add_textbox(s, Inches(0.5), Inches(3.15), Inches(6.2), Inches(0.35),
            "2️⃣ 重连失效 — 异步时序陷阱",
            font_size=BODY_SM, bold=True, color=DARK)
fix2 = """❌ newWebSocket() 后立即 send → 静默丢弃
✅ 在 onOpen 回调内 send，确保 ws OPEN"""
add_code_block(s, Inches(0.5), Inches(3.5), Inches(6.2), Inches(0.85),
                fix2, font_size=10)

add_textbox(s, Inches(0.5), Inches(4.45), Inches(6.2), Inches(0.35),
            "3️⃣ 两端结算不一致 — 统一公式",
            font_size=BODY_SM, bold=True, color=DARK)
fix3 = """playerScores 实时追踪每人已收
赢方 = 赢方已收 + 输方未走完(已收+手牌)
✓ 15 个验证用例全部通过"""
add_code_block(s, Inches(0.5), Inches(4.8), Inches(6.2), Inches(1.1),
                fix3, font_size=10)

# 右半：3 个部署 / 平台陷阱
add_textbox(s, Inches(7.0), Inches(1.05), Inches(6.0), Inches(0.4),
            "🚀 部署 / 平台层", font_size=BODY_LG, bold=True, color=ORANGE)

add_textbox(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(0.35),
            "4️⃣ Web 中文豆腐块（CJK）",
            font_size=BODY_SM, bold=True, color=DARK)
fix4 = """根因：浏览器无系统中文字体
方案：打包 ~3MB GB2312 子集字体
教训：跨端要测真实终端环境"""
add_code_block(s, Inches(7.0), Inches(1.85), Inches(6.0), Inches(1.2),
                fix4, font_size=10)

add_textbox(s, Inches(7.0), Inches(3.15), Inches(6.0), Inches(0.35),
            "5️⃣ 双层防火墙 — 部署对称陷阱",
            font_size=BODY_SM, bold=True, color=DARK)
fix5 = """ufw 开了 ≠ 公网通：腾讯云安全组也要开
教训：私有云常有"两层防火墙"，必须都查"""
add_code_block(s, Inches(7.0), Inches(3.5), Inches(6.0), Inches(0.85),
                fix5, font_size=10)

add_textbox(s, Inches(7.0), Inches(4.45), Inches(6.0), Inches(0.35),
            "6️⃣ Android URL 漂移 — 拓扑变更盲区",
            font_size=BODY_SM, bold=True, color=DARK)
fix6 = """PR #41 改 Caddy 反代后，Web 改了 URL
但 Android 硬编码 :8080 没改 → 完全连不上
教训：拓扑变更必须 grep 所有客户端 URL"""
add_code_block(s, Inches(7.0), Inches(4.8), Inches(6.0), Inches(1.1),
                fix6, font_size=10)

add_callout(s, Inches(0.5), Inches(6.05), Inches(12.5), Inches(1.0),
            "", bg=LIGHT_BG, border=PRIMARY)
add_textbox(s, Inches(0.8), Inches(6.2), Inches(12.0), Inches(0.4),
            "💡 共同规律：游戏逻辑层错在「假设」，部署层错在「跨端 / 跨拓扑同步」",
            font_size=BODY_MD, bold=True, color=PRIMARY)
add_textbox(s, Inches(0.8), Inches(6.6), Inches(12.0), Inches(0.4),
            "防御：(1) 关键路径强制 TDD；(2) 跨端共享 KMP 模块；(3) 拓扑变更 grep 所有 client 配置",
            font_size=BODY_SM, color=DARK)

# =================================================================
# Slide 13: Lessons + Action Items + Closing
# =================================================================
s = add_slide()
add_header(s, "八、经验总结 & 后续行动")

add_textbox(s, Inches(0.5), Inches(1.05), Inches(6.2), Inches(0.4),
            "🎯 核心经验",
            font_size=BODY_LG, bold=True, color=PRIMARY)
lessons = [
    "1. 「修了又坏」的根因：只修表层，没往深挖（需要分层根因审查）",
    "2. 跨端共享逻辑必须用 KMP 强制（已实现 :shared 模块）",
    "3. AI 静态扫描 + 人工真机 + Codex bot 不可互相替代（4 视角全跑）",
    "4. 协议 / 环境 / 部署拓扑变更必须人工打通（grep 所有客户端）",
    "5. 大特性 Phase 分段提交（协议 → 主客户端 → 次客户端）",
    "6. delay() 之后必须重检所有依赖项，不只 currentPlayerIndex",
]
y = 1.5
for l in lessons:
    add_textbox(s, Inches(0.6), Inches(y), Inches(6.2), Inches(0.35),
                l, font_size=BODY_SM, color=DARK)
    y += 0.4

add_textbox(s, Inches(7.0), Inches(1.05), Inches(6.0), Inches(0.4),
            "🚀 后续建议行动（已部分落地）",
            font_size=BODY_LG, bold=True, color=PRIMARY)
actions = [
    ["✅ 已实现", ":shared KMP 共享模块（PR #35）"],
    ["✅ 已实现", "PROTOCOL_VERSION + 握手（PR-H3）"],
    ["✅ 已实现", "harness L0-L4（hooks/playbook/regression）"],
    ["✅ 已实现", "tdd-gate CI 硬关（PR-H2）"],
    ["✅ 已实现", "pr-reviewer subagent（PR-H5）"],
    ["⚪ 待规划", "force-advance 监控告警 / 弱网 e2e"],
]
add_table(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(2.6),
            ["状态", "行动"], actions, font_size=BODY_SM,
            col_widths=[Inches(1.5), Inches(4.5)])

add_callout(s, Inches(0.5), Inches(4.4), Inches(12.4), Inches(2.3),
            "", bg=LIGHT_BG, border=PRIMARY)
add_textbox(s, Inches(0.8), Inches(4.55), Inches(12.0), Inches(0.5),
            "🏆 全程交付成果",
            font_size=BODY_LG, bold=True, color=PRIMARY)
add_textbox(s, Inches(0.8), Inches(5.05), Inches(12.0), Inches(0.4),
            "·  全程 54 PR / ~170 commit  ·  有效开发 18 天  ·  发现 ~128 个问题",
            font_size=BODY_SM, color=DARK)
add_textbox(s, Inches(0.8), Inches(5.45), Inches(12.0), Inches(0.4),
            "·  Android + Web 双端  ·  KMP 共享模块（编译期保证一致）  ·  Caddy 自托管自动部署",
            font_size=BODY_SM, color=DARK)
add_textbox(s, Inches(0.8), Inches(5.85), Inches(12.0), Inches(0.4),
            "·  游戏永不卡死 ✓  ·  重连无缝恢复 ✓  ·  两端结算一致 ✓  ·  PROTOCOL_VERSION = 3 ✓",
            font_size=BODY_SM, color=GREEN, bold=True)
add_textbox(s, Inches(0.8), Inches(6.3), Inches(12.0), Inches(0.4),
            "💡 核心结论：4 视角覆盖率 + harness 跨会话沉淀 = 教训不浪费",
            font_size=BODY_MD, bold=True, color=PRIMARY)

add_textbox(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
            "谢谢  ·  问题 & 讨论",
            font_size=BODY_MD, bold=True, color=GRAY,
            align=PP_ALIGN.CENTER)

# =================================================================
# Slide 14: 九、harness 跨会话经验 (1/3) — Phase 分段 + wasmJs
# =================================================================
s = add_slide()
add_header(s, "九、harness 跨会话经验（1/3）")

add_textbox(s, Inches(0.5), Inches(1.05), Inches(6.2), Inches(0.4),
            "📦 大特性的 Phase 分段提交",
            font_size=BODY_LG, bold=True, color=PRIMARY)

phase_rows = [
    ["Phase 1", "协议 + 服务端 + 单测", "底层稳了再动客户端"],
    ["Phase 2", "主客户端 (Android)", "一份吃通验证 server"],
    ["Phase 3", "次客户端 (Web) + 跨端测", "最后补齐"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(6.2), Inches(2.0),
            ["Phase", "内容", "价值"], phase_rows, font_size=BODY_SM,
            col_widths=[Inches(1.0), Inches(2.6), Inches(2.6)])

add_textbox(s, Inches(0.5), Inches(3.65), Inches(6.2), Inches(1.2),
            "适用：跨协议 + 服务端 + 双客户端 > 200 行 / > 5 文件\n"
            "每 Phase 内还按编译单元切（Android / Web 拆 commit）\n"
            "PR #53 G34-G38（5 特性, ~700 行）3 phase 实战验证",
            font_size=BODY_SM, color=DARK)

add_textbox(s, Inches(7.0), Inches(1.05), Inches(6.0), Inches(0.4),
            "⚠️  wasmJs 编译器隐藏陷阱",
            font_size=BODY_LG, bold=True, color=RED)

add_code_block(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(1.4),
               "Backend Internal error:\n"
               "  Exception during psi2ir\n"
               "Caused by: NullPointerException",
               font_size=BODY_SM)

add_textbox(s, Inches(7.0), Inches(3.05), Inches(6.0), Inches(2.3),
            "根因：可空 lambda + Compose smart-cast；\n"
            "vm::method KFunction → (() -> Unit)? 自动转换\n\n"
            "修法：\n"
            "• 可空 lambda → local val 固化非空\n"
            "• 函数引用 → 显式 { vm.foo() } lambda\n\n"
            "沙箱拉不到 wasmJs 编译器；写完 Web UI 必须 push CI",
            font_size=BODY_SM, color=DARK)

add_callout(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.7),
            "", bg=LIGHT_BG, border=PRIMARY)
add_textbox(s, Inches(0.8), Inches(5.15), Inches(12.0), Inches(0.4),
            "💡 实战教训：",
            font_size=BODY_MD, bold=True, color=PRIMARY)
add_textbox(s, Inches(0.8), Inches(5.5), Inches(12.0), Inches(1.2),
            "•  jvmTest 过 ≠ wasmJs 过 —— Web UI 改完必须跑 CI 才算\n"
            "•  Phase 内还要按编译单元切，平台特定编译错误更早暴露\n"
            "•  CI 红时先看 e: 开头的硬错误行，忽略 ~50 行 w: 警告噪音",
            font_size=BODY_SM, color=DARK)

# =================================================================
# Slide 15: 九、harness 跨会话经验 (2/3) — Codex 互补 + delay 重检
# =================================================================
s = add_slide()
add_header(s, "九、harness 跨会话经验（2/3）")

add_textbox(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
            "🔍 Codex bot 与 Claude /review-pr 盲区互补（PR #53 实证）",
            font_size=BODY_LG, bold=True, color=PRIMARY)

reviewer_rows = [
    ["Claude /review-pr", "docs 误写 + 引用不存在函数（**跨文件契约**）", "delay() 后 race"],
    ["Codex bot", "delay() 后没重检 isAISubstitute 的 race（**语句级边界**）", "文档与代码语义对齐"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
            ["Reviewer", "找到", "漏掉"], reviewer_rows, font_size=BODY_SM,
            col_widths=[Inches(2.5), Inches(6.3), Inches(3.5)])

add_callout(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.85),
            "", bg=LIGHT_BG, border=PRIMARY)
add_textbox(s, Inches(0.8), Inches(3.3), Inches(12.0), Inches(0.6),
            "盲区不重叠 → 单独跑一个至少漏一类问题。Codex + Claude 必须都跑。",
            font_size=BODY_MD, bold=True, color=DARK)

add_textbox(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.4),
            "⏰ delay() 后状态过期 — race 反请（regressions #11）",
            font_size=BODY_LG, bold=True, color=RED)

add_code_block(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.4),
               "delay(effectiveAiDelayMs(...))     // 玩家此时取消了托管\n"
               "mutexFor(room).withLock {\n"
               "    if (state.currentPlayerIndex != playerIndex) return  // 只检查了这一项\n"
               "    decideAIAction(...)             // isAISubstitute 已变 → 不该再代打\n"
               "}", font_size=BODY_SM)

add_textbox(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.7),
            "审查清单：醒来后必须重检「延迟前所有依赖项」，不只检查最显眼的那个\n"
            "修法：抽 internal fun shouldYieldToHumanPlayer(...) 谓词 + 同 commit 加测",
            font_size=BODY_SM, color=DARK)

# =================================================================
# Slide 16: 九、harness 跨会话经验 (3/3) — L0-L4 体系 + 4 教训
# =================================================================
s = add_slide()
add_header(s, "九、harness 跨会话经验（3/3）— L0–L4 体系")

add_textbox(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
            "🏗️  Harness 5 层防御体系（PR-H1..H5 + #35 web 重构 实战搭建）",
            font_size=BODY_LG, bold=True, color=PRIMARY)

# 复用 build_html.py 的 L0-L4 SVG（与渲染版 HTML 视觉一致）
# 左半边 SVG，右半边经验表 — 2 栏布局避免上下挤压
embed_svg(s, svg_harness_l0_l4(), Inches(0.4), Inches(1.5), width=Inches(6.6))

add_textbox(s, Inches(7.2), Inches(1.05), Inches(6.0), Inches(0.4),
            "🎯 跨会话经验的核心原则",
            font_size=BODY_LG, bold=True, color=GREEN)

other_lessons = [
    ["同 commit *Test.kt 配对", "tdd-gate 一次没误报、一次没漏；hook 弹\"TDD 提醒\"时不必另开 commit"],
    ["分支 vs PR 一一对应", "PR merge 后开新分支；老分支再 push 不会自动出 PR（PR #52→#53 实战）"],
    ["文档单一真相 grep 验证", "新 doc 自称\"权威\"前必 grep 验证 anchor 函数名实际存在"],
    ["速度档位用 3 档预设", "slider 看似灵活；实际多数场景 3 档够用 + 测试矩阵更小"],
]
add_table(s, Inches(7.2), Inches(1.5), Inches(5.9), Inches(4.6),
            ["经验", "说明"], other_lessons, font_size=BODY_SM,
            col_widths=[Inches(2.0), Inches(3.9)])

add_callout(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.55),
            "💎 核心：harness 让教训跨 session 沉淀，新人 / 新 AI 不需要每次从头踩坑",
            bg=PRIMARY, border=PRIMARY,
            font_size=BODY_MD, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# =================================================================
total = len(prs.slides)
for i, slide in enumerate(prs.slides, start=1):
    if i == 1:
        continue
    add_page_number(slide, i, total)

output = "/home/user/AndroidAPP/docs/dev_summary.pptx"
prs.save(output)
print(f"✓ Saved {total} slides to {output}")
